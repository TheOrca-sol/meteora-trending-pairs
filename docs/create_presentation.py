#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Meteora Trending Pairs
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Sodexo Live / StayManager Colors
PRIMARY = RGBColor(228, 0, 124)  # #E4007C
SECONDARY = RGBColor(124, 36, 131)  # #7C2483
GOLD = RGBColor(255, 215, 0)  # #FFD700
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(26, 26, 46)
LIGHT_GRAY = RGBColor(245, 245, 245)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Hero Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = PRIMARY

    # Badge
    badge = slide1.shapes.add_textbox(Inches(3.5), Inches(1.5), Inches(3), Inches(0.5))
    badge_frame = badge.text_frame
    badge_frame.text = "✨ REVOLUTIONARY PLATFORM ✨"
    badge_p = badge_frame.paragraphs[0]
    badge_p.font.size = Pt(14)
    badge_p.font.bold = True
    badge_p.font.color.rgb = DARK
    badge_p.alignment = PP_ALIGN.CENTER
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD

    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Meteora Trending Pairs"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(2), Inches(3.8), Inches(6), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Real-Time DLMM Analytics Platform"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = WHITE
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide1.shapes.add_textbox(Inches(2.5), Inches(4.6), Inches(5), Inches(0.5))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Track 4,000+ pools • Built-in security • Automated alerts"
    tagline_p = tagline_frame.paragraphs[0]
    tagline_p.font.size = Pt(16)
    tagline_p.font.color.rgb = WHITE
    tagline_p.alignment = PP_ALIGN.CENTER

    # Stats
    stats_y = 5.5
    stats = [("4K+", "Pools Monitored"), ("60s", "Refresh Rate"), ("24/7", "Automated Alerts")]
    for i, (number, label) in enumerate(stats):
        x = 1.5 + (i * 2.5)
        stat_box = slide1.shapes.add_textbox(Inches(x), Inches(stats_y), Inches(2), Inches(1))
        stat_frame = stat_box.text_frame
        stat_frame.text = f"{number}\n{label}"
        stat_p = stat_frame.paragraphs[0]
        stat_p.font.size = Pt(32)
        stat_p.font.bold = True
        stat_p.font.color.rgb = WHITE
        stat_p.alignment = PP_ALIGN.CENTER
        stat_frame.paragraphs[1].font.size = Pt(12)
        stat_frame.paragraphs[1].font.color.rgb = WHITE
        stat_frame.paragraphs[1].alignment = PP_ALIGN.CENTER

    # Slide 2: The Challenge
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = LIGHT_GRAY

    title = slide2.shapes.title
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "The Challenge"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY

    # Gradient box
    gradient_box = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.2))
    gradient_box.fill.solid()
    gradient_box.fill.fore_color.rgb = PRIMARY
    gradient_frame = gradient_box.text_frame
    gradient_frame.text = "Traders & LPs Face Critical Obstacles\n\nFinding profitable opportunities in 4,000+ Meteora DLMM pools is overwhelming"
    for p in gradient_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    gradient_frame.paragraphs[0].font.size = Pt(24)
    gradient_frame.paragraphs[0].font.bold = True
    gradient_frame.paragraphs[1].font.size = Pt(16)

    # Feature cards
    challenges = [
        ("1", "Information Overload", "Too many pools, not enough context. Manual research takes hours."),
        ("2", "Security Risks", "Rug pulls and scams are rampant. No unified security analysis."),
        ("3", "Missed Opportunities", "High-APR pools trend for minutes. By the time you notice, it's too late."),
        ("4", "Data Fragmentation", "Information scattered across 7+ sources. No single dashboard.")
    ]

    for i, (num, title_text, desc) in enumerate(challenges):
        row = i // 2
        col = i % 2
        x = 0.8 + (col * 4.5)
        y = 3.2 + (row * 1.8)

        card = slide2.shapes.add_textbox(Inches(x), Inches(y), Inches(4), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card_frame = card.text_frame
        card_frame.text = f"{num}\n{title_text}\n{desc}"
        card_frame.paragraphs[0].font.size = Pt(28)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = PRIMARY
        card_frame.paragraphs[1].font.size = Pt(18)
        card_frame.paragraphs[1].font.bold = True
        card_frame.paragraphs[1].font.color.rgb = DARK
        card_frame.paragraphs[2].font.size = Pt(12)
        card_frame.paragraphs[2].font.color.rgb = RGBColor(102, 102, 102)

    # Slide 3: Our Solution
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame3 = title_box3.text_frame
    title_frame3.text = "Our Solution"
    title_p3 = title_frame3.paragraphs[0]
    title_p3.font.size = Pt(48)
    title_p3.font.bold = True
    title_p3.font.color.rgb = PRIMARY

    # Gradient box
    gradient_box3 = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.2))
    gradient_box3.fill.solid()
    gradient_box3.fill.fore_color.rgb = PRIMARY
    gradient_frame3 = gradient_box3.text_frame
    gradient_frame3.text = "A Single Platform for Complete DLMM Intelligence\n\nWe aggregate data from 7+ sources, analyze security, and deliver real-time alerts"
    for p in gradient_frame3.paragraphs:
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    gradient_frame3.paragraphs[0].font.size = Pt(24)
    gradient_frame3.paragraphs[0].font.bold = True
    gradient_frame3.paragraphs[1].font.size = Pt(16)

    # Stats
    stats3 = [
        ("4K+", "Pools Monitored", "Every Meteora DLMM pool, updated every 60 seconds"),
        ("7+", "Data Sources", "Meteora, DexScreener, Jupiter, RugCheck, Helius & more"),
        ("24/7", "Automated Alerts", "Telegram notifications for trending opportunities")
    ]

    for i, (number, title_text, desc) in enumerate(stats3):
        x = 0.8 + (i * 3)
        card = slide3.shapes.add_textbox(Inches(x), Inches(3.2), Inches(2.8), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card_frame = card.text_frame
        card_frame.text = f"{number}\n{title_text}\n{desc}"
        card_frame.paragraphs[0].font.size = Pt(36)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = PRIMARY
        card_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        card_frame.paragraphs[1].font.size = Pt(16)
        card_frame.paragraphs[1].font.bold = True
        card_frame.paragraphs[1].font.color.rgb = DARK
        card_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
        card_frame.paragraphs[2].font.size = Pt(11)
        card_frame.paragraphs[2].font.color.rgb = RGBColor(102, 102, 102)
        card_frame.paragraphs[2].alignment = PP_ALIGN.CENTER

    # Slide 4: Core Features
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame4 = title_box4.text_frame
    title_frame4.text = "Core Features"
    title_p4 = title_frame4.paragraphs[0]
    title_p4.font.size = Pt(48)
    title_p4.font.bold = True
    title_p4.font.color.rgb = PRIMARY

    features = [
        ("📊", "Real-Time Analytics", "Monitor 30m fee rates, 24h volume, APR, TVL with auto-refresh"),
        ("🔍", "Advanced Filtering", "Find what you need with filters for APR, volume, TVL, fees"),
        ("🛡️", "Security Analysis", "Integrated RugCheck scans, holder distribution, authority checks"),
        ("📱", "Telegram Alerts", "Instant notifications when pools match your criteria"),
        ("⚡", "Degen Mode", "High-frequency monitoring (1-60 min) for active traders"),
        ("📈", "Deep Analytics", "Complete pool data: holder info, BubbleMaps, volume trends")
    ]

    for i, (icon, title_text, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = 0.8 + (col * 4.5)
        y = 1.8 + (row * 1.6)

        card = slide4.shapes.add_textbox(Inches(x), Inches(y), Inches(4), Inches(1.3))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card_frame = card.text_frame
        card_frame.text = f"{icon} {title_text}\n{desc}"
        card_frame.paragraphs[0].font.size = Pt(18)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = DARK
        card_frame.paragraphs[1].font.size = Pt(12)
        card_frame.paragraphs[1].font.color.rgb = RGBColor(102, 102, 102)

    # Continue with remaining slides...
    # Slide 5-14 would follow similar pattern

    # Slide 14: Closing Hero
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    slide14.background.fill.solid()
    slide14.background.fill.fore_color.rgb = PRIMARY

    badge14 = slide14.shapes.add_textbox(Inches(3), Inches(1.5), Inches(4), Inches(0.5))
    badge14_frame = badge14.text_frame
    badge14_frame.text = "✨ THE FUTURE OF DLMM TRADING ✨"
    badge14_p = badge14_frame.paragraphs[0]
    badge14_p.font.size = Pt(14)
    badge14_p.font.bold = True
    badge14_p.font.color.rgb = DARK
    badge14_p.alignment = PP_ALIGN.CENTER
    badge14.fill.solid()
    badge14.fill.fore_color.rgb = GOLD

    title_box14 = slide14.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame14 = title_box14.text_frame
    title_frame14.text = "Ready to Trade\nSmarter?"
    title_p14 = title_frame14.paragraphs[0]
    title_p14.font.size = Pt(60)
    title_p14.font.bold = True
    title_p14.font.color.rgb = WHITE
    title_p14.alignment = PP_ALIGN.CENTER

    subtitle_box14 = slide14.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(0.5))
    subtitle_frame14 = subtitle_box14.text_frame
    subtitle_frame14.text = "Join us in revolutionizing DLMM analytics on Solana"
    subtitle_p14 = subtitle_frame14.paragraphs[0]
    subtitle_p14.font.size = Pt(20)
    subtitle_p14.font.color.rgb = WHITE
    subtitle_p14.alignment = PP_ALIGN.CENTER

    thanks_box = slide14.shapes.add_textbox(Inches(2.5), Inches(5.5), Inches(5), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank you for your time!\n\nQuestions? Let's discuss how we can work together."
    for p in thanks_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    thanks_frame.paragraphs[0].font.size = Pt(20)
    thanks_frame.paragraphs[0].font.bold = True
    thanks_frame.paragraphs[1].font.size = Pt(16)

    return prs

if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/younes/Documents/work/meteora-trending-pairs/meteora-trending-pairs/docs/Meteora_Trending_Pairs_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint presentation created: {output_path}")
