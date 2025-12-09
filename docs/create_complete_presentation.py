#!/usr/bin/env python3
"""
Generate Complete PowerPoint presentation for Meteora Trending Pairs
With all 14 slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Colors
PRIMARY = RGBColor(228, 0, 124)  # #E4007C
SECONDARY = RGBColor(124, 36, 131)  # #7C2483
GOLD = RGBColor(255, 215, 0)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(26, 26, 46)
LIGHT_GRAY = RGBColor(250, 250, 250)
GRAY = RGBColor(102, 102, 102)

def add_title_slide(prs, badge_text, title_text, subtitle_text, tagline_text):
    """Add a hero title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PRIMARY

    # Badge
    badge = slide.shapes.add_textbox(Inches(3), Inches(1.2), Inches(4), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    tf = badge.text_frame
    tf.text = badge_text
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(2), Inches(3.6), Inches(6), Inches(0.5))
    tf = subtitle.text_frame
    tf.text = subtitle_text
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline = slide.shapes.add_textbox(Inches(2.5), Inches(4.3), Inches(5), Inches(0.4))
    tf = tagline.text_frame
    tf.text = tagline_text
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Stats
    stats = [("4K+", "Pools"), ("60s", "Refresh"), ("24/7", "Alerts")]
    for i, (num, label) in enumerate(stats):
        x = 2 + (i * 2)
        box = slide.shapes.add_textbox(Inches(x), Inches(5.2), Inches(1.8), Inches(1))
        tf = box.text_frame
        tf.text = f"{num}\n{label}"
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[1].font.size = Pt(12)
        tf.paragraphs[1].font.color.rgb = WHITE
        tf.paragraphs[1].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title_text):
    """Add a content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Underline
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(1.2), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()

    return slide

def add_gradient_box(slide, y, title, subtitle):
    """Add a gradient box with text"""
    box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = PRIMARY
    tf = box.text_frame
    tf.text = f"{title}\n{subtitle}"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(14)
    tf.paragraphs[1].font.color.rgb = WHITE
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER

def add_feature_cards(slide, features, start_y):
    """Add feature cards in 2x2 grid"""
    for i, (icon, title, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = 0.8 + (col * 4.5)
        y = start_y + (row * 1.6)

        card = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4), Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        tf = card.text_frame
        tf.text = f"{icon} {title}\n{desc}"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = DARK
        tf.paragraphs[1].font.size = Pt(11)
        tf.paragraphs[1].font.color.rgb = GRAY

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Hero
add_title_slide(prs, "✨ REVOLUTIONARY PLATFORM ✨",
                "Meteora Trending Pairs",
                "Real-Time DLMM Analytics Platform",
                "Track 4,000+ pools • Built-in security • Automated alerts")

# Slide 2: Challenge
slide2 = add_content_slide(prs, "The Challenge")
add_gradient_box(slide2, 1.4, "Traders & LPs Face Critical Obstacles",
                 "Finding profitable opportunities in 4,000+ Meteora DLMM pools is overwhelming")
features2 = [
    ("1", "Information Overload", "Too many pools, not enough context"),
    ("2", "Security Risks", "Rug pulls and scams are rampant"),
    ("3", "Missed Opportunities", "High-APR pools trend for minutes"),
    ("4", "Data Fragmentation", "Information scattered across 7+ sources")
]
add_feature_cards(slide2, features2, 2.8)

# Slide 3: Solution
slide3 = add_content_slide(prs, "Our Solution")
add_gradient_box(slide3, 1.4, "A Single Platform for Complete DLMM Intelligence",
                 "We aggregate data from 7+ sources and deliver real-time alerts")
stats3 = [("4K+", "Pools Monitored"), ("7+", "Data Sources"), ("24/7", "Automated Alerts")]
for i, (num, label) in enumerate(stats3):
    x = 1 + (i * 2.8)
    card = slide3.shapes.add_textbox(Inches(x), Inches(2.9), Inches(2.6), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    tf = card.text_frame
    tf.text = f"{num}\n{label}"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = PRIMARY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(14)
    tf.paragraphs[1].font.bold = True
    tf.paragraphs[1].font.color.rgb = DARK
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER

# Slide 4: Core Features
slide4 = add_content_slide(prs, "Core Features")
features4 = [
    ("📊", "Real-Time Analytics", "Monitor 30m fee rates, 24h volume, APR, TVL"),
    ("🔍", "Advanced Filtering", "Find what you need with filters"),
    ("🛡️", "Security Analysis", "RugCheck scans, holder distribution"),
    ("📱", "Telegram Alerts", "Instant notifications"),
    ("⚡", "Degen Mode", "High-frequency monitoring"),
    ("📈", "Deep Analytics", "Complete pool data")
]
add_feature_cards(slide4, features4, 1.5)

# Slide 5: How It Works
slide5 = add_content_slide(prs, "How It Works")
steps = [
    ("1", "Data Aggregation", "Fetch pool data from multiple sources"),
    ("2", "Smart Analysis", "Process 4,000+ pools and run security checks"),
    ("3", "Intelligent Filtering", "Show matching opportunities"),
    ("4", "Real-Time Alerts", "Telegram notifications with trading links")
]
add_feature_cards(slide5, steps, 1.5)

# Slide 6: Tech Stack
slide6 = add_content_slide(prs, "Technology Stack")
tech = ["React 18", "Material-UI", "Flask", "Node.js", "PostgreSQL", "Supabase",
        "Solana Web3.js", "APScheduler", "Telegram Bot", "DexScreener", "Jupiter", "RugCheck"]
for i, tech_name in enumerate(tech):
    row = i // 3
    col = i % 3
    x = 1.2 + (col * 2.8)
    y = 1.8 + (row * 0.8)
    box = slide6.shapes.add_textbox(Inches(x), Inches(y), Inches(2.4), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    tf = box.text_frame
    tf.text = tech_name
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

# Slide 7: Competitive Advantages
slide7 = add_content_slide(prs, "Why We're Different")
features7 = [
    ("🎯", "100% DLMM Focused", "Built exclusively for Meteora"),
    ("🔒", "Security-First", "Integrated security analysis"),
    ("📲", "Mobile-First Alerts", "Always connected via Telegram"),
    ("⚙️", "Highly Customizable", "Adapts to your strategy")
]
add_feature_cards(slide7, features7, 1.5)

# Slide 8: Roadmap
slide8 = add_content_slide(prs, "Product Roadmap")
roadmap = [
    ("Q1 2025", ["Automated trading", "Portfolio tracking", "Backtesting", "Watchlist"]),
    ("Q2 2025", ["Custom alerts", "Multi-DEX support", "Mobile apps", "API access"]),
    ("Q3 2025", ["Social features", "AI-powered scoring", "Cross-chain", "Enterprise"])
]
for i, (quarter, items) in enumerate(roadmap):
    x = 1 + (i * 2.8)
    card = slide8.shapes.add_textbox(Inches(x), Inches(1.8), Inches(2.4), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    tf = card.text_frame
    tf.text = f"{quarter}\n" + "\n".join(f"• {item}" for item in items)
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = PRIMARY
    for j in range(1, len(tf.paragraphs)):
        tf.paragraphs[j].font.size = Pt(10)
        tf.paragraphs[j].font.color.rgb = GRAY

# Slide 9: Business Model
slide9 = add_content_slide(prs, "Business Model")
tiers = [
    ("Free Tier", "Core analytics\nBasic monitoring\nSecurity analysis\nAlways free"),
    ("Premium", "$29/month\n1-min monitoring\nUnlimited alerts\nAPI access"),
    ("Enterprise", "Custom pricing\nWhite-label\nDedicated infra\nSLA guarantees")
]
for i, (tier, desc) in enumerate(tiers):
    x = 1 + (i * 2.8)
    card = slide9.shapes.add_textbox(Inches(x), Inches(1.8), Inches(2.4), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    tf = card.text_frame
    tf.text = f"{tier}\n{desc}"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK
    for j in range(1, len(tf.paragraphs)):
        tf.paragraphs[j].font.size = Pt(11)
        tf.paragraphs[j].font.color.rgb = GRAY

# Slide 10: Market Opportunity
slide10 = add_content_slide(prs, "Market Opportunity")
features10 = [
    ("💰", "Growing DeFi Market", "Solana DeFi TVL exceeds $8B"),
    ("👥", "Target Audience", "Traders, LPs, protocols, trading firms"),
    ("📈", "Scalable Model", "Start Meteora, expand multi-chain"),
    ("🚀", "First-Mover", "No dedicated DLMM analytics platform")
]
add_feature_cards(slide10, features10, 1.5)

# Slide 11: Traction
slide11 = add_content_slide(prs, "Current Traction")
stats11 = [("100%", "Feature Complete"), ("4K+", "Pools Tracked"), ("60s", "Refresh Rate")]
for i, (num, label) in enumerate(stats11):
    x = 1.5 + (i * 2.5)
    card = slide11.shapes.add_textbox(Inches(x), Inches(1.8), Inches(2.2), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    tf = card.text_frame
    tf.text = f"{num}\n{label}"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = PRIMARY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(14)
    tf.paragraphs[1].font.bold = True
    tf.paragraphs[1].font.color.rgb = DARK
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER
add_gradient_box(slide11, 3.8, "Ready to Launch",
                 "Platform is fully functional and tested")

# Slide 12: Live Demo
slide12 = add_content_slide(prs, "Live Demo")
add_gradient_box(slide12, 2.5, "Experience the Platform Live",
                 "See real-time data from 4,000+ pools")
cta = slide12.shapes.add_textbox(Inches(3.5), Inches(4.2), Inches(3), Inches(0.6))
cta.fill.solid()
cta.fill.fore_color.rgb = GOLD
tf = cta.text_frame
tf.text = "Open Dashboard →"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER

# Slide 13: Contact
slide13 = add_content_slide(prs, "Let's Connect")
contacts = [("🌐", "Live Platform"), ("📧", "Email"), ("🐦", "Twitter"), ("💬", "Telegram")]
for i, (icon, label) in enumerate(contacts):
    row = i // 2
    col = i % 2
    x = 1.5 + (col * 3.8)
    y = 1.8 + (row * 1.5)
    card = slide13.shapes.add_textbox(Inches(x), Inches(y), Inches(3.4), Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    tf = card.text_frame
    tf.text = f"{icon} {label}\n[Your Info]"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(12)
    tf.paragraphs[1].font.color.rgb = PRIMARY
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER
add_gradient_box(slide13, 4.5, "We're Seeking",
                 "Strategic Partners • Advisors • Beta Users • Feedback")

# Slide 14: Closing Hero
add_title_slide(prs, "✨ THE FUTURE OF DLMM TRADING ✨",
                "Ready to Trade Smarter?",
                "Join us in revolutionizing DLMM analytics",
                "Thank you for your time!")

# Save
output_path = "/home/younes/Documents/work/meteora-trending-pairs/meteora-trending-pairs/docs/Meteora_Trending_Pairs_Presentation.pptx"
prs.save(output_path)
print(f"✅ PowerPoint presentation created successfully!")
print(f"📁 Location: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
