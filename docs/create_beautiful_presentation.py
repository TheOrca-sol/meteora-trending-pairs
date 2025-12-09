#!/usr/bin/env python3
"""
Generate Beautiful PowerPoint presentation for Meteora Trending Pairs
Inspired by StayManager design with decorative shapes and modern styling
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Sodexo Live / StayManager Colors
PRIMARY = RGBColor(228, 0, 124)  # #E4007C
SECONDARY = RGBColor(124, 36, 131)  # #7C2483
GOLD = RGBColor(255, 215, 0)  # #FFD700
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(26, 26, 46)
LIGHT_GRAY = RGBColor(245, 245, 245)
GRAY = RGBColor(102, 102, 102)
LIGHT_PRIMARY = RGBColor(255, 182, 218)  # Lighter pink for shapes

def add_decorative_shapes(slide, color1, color2):
    """Add floating decorative circles to hero slides like StayManager"""
    # Large circle - top right
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7.5), Inches(-1), Inches(4), Inches(4)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = color2
    circle1.fill.transparency = 0.3
    circle1.line.fill.background()

    # Medium circle - bottom left
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-1.5), Inches(5), Inches(3.5), Inches(3.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = LIGHT_PRIMARY
    circle2.fill.transparency = 0.4
    circle2.line.fill.background()

    # Small circle - middle right
    circle3 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(8), Inches(3), Inches(2.5), Inches(2.5)
    )
    circle3.fill.solid()
    circle3.fill.fore_color.rgb = color1
    circle3.fill.transparency = 0.35
    circle3.line.fill.background()

def add_rounded_card_with_shadow(slide, x, y, width, height):
    """Add a card with shadow effect"""
    # Shadow (slightly offset)
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.03), Inches(y + 0.03), Inches(width), Inches(height)
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(200, 200, 200)
    shadow.fill.transparency = 0.5
    shadow.line.fill.background()

    # Main card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()

    return card

def add_card_with_gradient_top(slide, x, y, width, height, gradient_color):
    """Add a card with gradient top bar like StayManager"""
    # Shadow
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.03), Inches(y + 0.03), Inches(width), Inches(height)
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(200, 200, 200)
    shadow.fill.transparency = 0.5
    shadow.line.fill.background()

    # Main card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()

    # Gradient top bar
    gradient_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(0.15)
    )
    gradient_bar.fill.solid()
    gradient_bar.fill.fore_color.rgb = gradient_color
    gradient_bar.line.fill.background()

    return card

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1: HERO TITLE =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = PRIMARY

    # Add decorative shapes
    add_decorative_shapes(slide1, WHITE, SECONDARY)

    # Golden badge with rounded corners
    badge = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.3), Inches(1.5), Inches(3.4), Inches(0.5)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()
    badge_frame = badge.text_frame
    badge_frame.text = "✨ REVOLUTIONARY PLATFORM ✨"
    badge_p = badge_frame.paragraphs[0]
    badge_p.font.size = Pt(14)
    badge_p.font.bold = True
    badge_p.font.color.rgb = DARK
    badge_p.alignment = PP_ALIGN.CENTER
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "Meteora Trending Pairs"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(64)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(2), Inches(3.7), Inches(6), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Real-Time DLMM Analytics Platform"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(26)
    subtitle_p.font.color.rgb = WHITE
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide1.shapes.add_textbox(Inches(2), Inches(4.4), Inches(6), Inches(0.4))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Track 4,000+ pools • Built-in security • Automated alerts"
    tagline_p = tagline_frame.paragraphs[0]
    tagline_p.font.size = Pt(15)
    tagline_p.font.color.rgb = RGBColor(255, 255, 255)
    tagline_p.alignment = PP_ALIGN.CENTER

    # Stats boxes with rounded corners
    stats_y = 5.3
    stats = [("4K+", "Pools"), ("60s", "Refresh"), ("24/7", "Alerts")]
    for i, (number, label) in enumerate(stats):
        x = 1.7 + (i * 2.2)
        stat_box = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(stats_y), Inches(2), Inches(0.9)
        )
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        stat_box.fill.transparency = 0.15
        stat_box.line.fill.background()

        stat_frame = stat_box.text_frame
        stat_frame.text = f"{number}\n{label}"
        stat_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        stat_frame.paragraphs[0].font.size = Pt(30)
        stat_frame.paragraphs[0].font.bold = True
        stat_frame.paragraphs[0].font.color.rgb = WHITE
        stat_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        stat_frame.paragraphs[1].font.size = Pt(11)
        stat_frame.paragraphs[1].font.color.rgb = WHITE
        stat_frame.paragraphs[1].alignment = PP_ALIGN.CENTER

    # ===== SLIDE 2: THE CHALLENGE =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = LIGHT_GRAY

    # Title
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "The Challenge"
    title_p2 = title_frame2.paragraphs[0]
    title_p2.font.size = Pt(52)
    title_p2.font.bold = True
    title_p2.font.color.rgb = PRIMARY

    # Gradient box with rounded corners
    gradient_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.35), Inches(8.4), Inches(1.2)
    )
    gradient_box.fill.solid()
    gradient_box.fill.fore_color.rgb = PRIMARY
    gradient_box.line.fill.background()
    gradient_frame = gradient_box.text_frame
    gradient_frame.text = "Traders & LPs Face Critical Obstacles\n\nFinding profitable opportunities in 4,000+ Meteora DLMM pools is overwhelming"
    gradient_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in gradient_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    gradient_frame.paragraphs[0].font.size = Pt(26)
    gradient_frame.paragraphs[0].font.bold = True
    gradient_frame.paragraphs[2].font.size = Pt(15)

    # Challenge cards
    challenges = [
        ("1", "Information Overload", "Too many pools, not enough context"),
        ("2", "Security Risks", "Rug pulls and scams are rampant"),
        ("3", "Missed Opportunities", "High-APR pools trend for minutes"),
        ("4", "Data Fragmentation", "Information scattered across 7+ sources")
    ]

    for i, (num, title_text, desc) in enumerate(challenges):
        row = i // 2
        col = i % 2
        x = 0.8 + (col * 4.4)
        y = 3 + (row * 1.7)

        card = add_card_with_gradient_top(slide2, x, y, 4.0, 1.4, PRIMARY)
        card_frame = card.text_frame
        card_frame.text = f"{num}\n{title_text}\n{desc}"
        card_frame.paragraphs[0].font.size = Pt(30)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = PRIMARY
        card_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        card_frame.paragraphs[1].font.size = Pt(17)
        card_frame.paragraphs[1].font.bold = True
        card_frame.paragraphs[1].font.color.rgb = DARK
        card_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
        card_frame.paragraphs[2].font.size = Pt(11)
        card_frame.paragraphs[2].font.color.rgb = GRAY
        card_frame.paragraphs[2].alignment = PP_ALIGN.CENTER

    # ===== SLIDE 3: OUR SOLUTION =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = LIGHT_GRAY

    # Title
    title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame3 = title_box3.text_frame
    title_frame3.text = "Our Solution"
    title_p3 = title_frame3.paragraphs[0]
    title_p3.font.size = Pt(52)
    title_p3.font.bold = True
    title_p3.font.color.rgb = PRIMARY

    # Gradient box
    gradient_box3 = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.35), Inches(8.4), Inches(1.2)
    )
    gradient_box3.fill.solid()
    gradient_box3.fill.fore_color.rgb = PRIMARY
    gradient_box3.line.fill.background()
    gradient_frame3 = gradient_box3.text_frame
    gradient_frame3.text = "A Single Platform for Complete DLMM Intelligence\n\nWe aggregate data from 7+ sources and deliver real-time alerts"
    gradient_frame3.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in gradient_frame3.paragraphs:
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    gradient_frame3.paragraphs[0].font.size = Pt(26)
    gradient_frame3.paragraphs[0].font.bold = True
    gradient_frame3.paragraphs[2].font.size = Pt(15)

    # Solution stats
    stats3 = [
        ("4K+", "Pools Monitored", "Every Meteora DLMM pool"),
        ("8+", "Data Sources", "Meteora, DexScreener, Jupiter, RugCheck & more"),
        ("24/7", "Automated Alerts", "Telegram notifications")
    ]

    for i, (number, title_text, desc) in enumerate(stats3):
        x = 0.9 + (i * 3)
        card = add_card_with_gradient_top(slide3, x, y=3, width=2.8, height=2.3, gradient_color=PRIMARY)
        card_frame = card.text_frame
        card_frame.text = f"{number}\n{title_text}\n{desc}"
        card_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        card_frame.paragraphs[0].font.size = Pt(38)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = PRIMARY
        card_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        card_frame.paragraphs[1].font.size = Pt(15)
        card_frame.paragraphs[1].font.bold = True
        card_frame.paragraphs[1].font.color.rgb = DARK
        card_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
        card_frame.paragraphs[2].font.size = Pt(11)
        card_frame.paragraphs[2].font.color.rgb = GRAY
        card_frame.paragraphs[2].alignment = PP_ALIGN.CENTER

    # ===== SLIDE 4: CORE FEATURES =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame4 = title_box4.text_frame
    title_frame4.text = "Core Features"
    title_p4 = title_frame4.paragraphs[0]
    title_p4.font.size = Pt(52)
    title_p4.font.bold = True
    title_p4.font.color.rgb = PRIMARY

    features = [
        ("📊", "Real-Time Analytics", "Monitor 30m fees, 24h volume, APR, TVL"),
        ("🔄", "Aggregated Liquidity", "Combined depth across all pools per token pair"),
        ("💰", "Capital Rotation", "24/7 monitoring with Telegram alerts for LPs"),
        ("🛡️", "Security Analysis", "RugCheck, holder distribution, authority checks"),
        ("🔍", "Advanced Filtering", "Find pools with custom filters & thresholds"),
        ("📱", "Telegram Alerts", "Instant notifications when opportunities appear")
    ]

    for i, (icon, title_text, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = 0.8 + (col * 4.4)
        y = 1.7 + (row * 1.6)

        card = add_card_with_gradient_top(slide4, x, y, 4.0, 1.35, SECONDARY)
        card_frame = card.text_frame
        card_frame.text = f"{icon}  {title_text}\n{desc}"
        card_frame.paragraphs[0].font.size = Pt(17)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = DARK
        card_frame.paragraphs[1].font.size = Pt(12)
        card_frame.paragraphs[1].font.color.rgb = GRAY

    # ===== SLIDE 5: AGGREGATED LIQUIDITY (NEW) =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame5 = title_box5.text_frame
    title_frame5.text = "Aggregated Liquidity"
    title_p5 = title_frame5.paragraphs[0]
    title_p5.font.size = Pt(52)
    title_p5.font.bold = True
    title_p5.font.color.rgb = PRIMARY

    # Explanation box
    explain_box = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.35), Inches(8), Inches(1.1)
    )
    explain_box.fill.solid()
    explain_box.fill.fore_color.rgb = PRIMARY
    explain_box.line.fill.background()
    explain_frame = explain_box.text_frame
    explain_frame.text = "Unified Market Depth View\n\nCombine liquidity from all DLMM pools with the same token pair to see true market depth"
    explain_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in explain_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    explain_frame.paragraphs[0].font.size = Pt(26)
    explain_frame.paragraphs[0].font.bold = True
    explain_frame.paragraphs[2].font.size = Pt(15)

    # Benefits
    benefits = [
        ("🎯", "True Market Depth", "See total liquidity across all configurations"),
        ("📊", "Better Trading Decisions", "Understand real slippage and depth"),
        ("🔍", "Multi-Pool Analysis", "Compare different bin steps side-by-side"),
        ("💡", "Unique Insight", "No other platform offers aggregated DLMM liquidity")
    ]

    for i, (icon, title_text, desc) in enumerate(benefits):
        row = i // 2
        col = i % 2
        x = 0.8 + (col * 4.4)
        y = 2.9 + (row * 1.6)

        card = add_card_with_gradient_top(slide5, x, y, 4.0, 1.35, GOLD)
        card_frame = card.text_frame
        card_frame.text = f"{icon}  {title_text}\n{desc}"
        card_frame.paragraphs[0].font.size = Pt(17)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = DARK
        card_frame.paragraphs[1].font.size = Pt(12)
        card_frame.paragraphs[1].font.color.rgb = GRAY

    # ===== SLIDE 6: HOW IT WORKS =====
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame6 = title_box6.text_frame
    title_frame6.text = "How It Works"
    title_p6 = title_frame6.paragraphs[0]
    title_p6.font.size = Pt(52)
    title_p6.font.bold = True
    title_p6.font.color.rgb = PRIMARY

    steps = [
        ("1", "Data Aggregation", "Fetch from 8+ sources (Meteora, DexScreener, Jupiter...)"),
        ("2", "Smart Analysis", "Process 4K+ pools & run security scans"),
        ("3", "Intelligent Filtering", "Apply filters & show opportunities"),
        ("4", "Real-Time Alerts", "Send Telegram notifications instantly")
    ]

    for i, (num, title_text, desc) in enumerate(steps):
        row = i // 2
        col = i % 2
        x = 0.8 + (col * 4.4)
        y = 1.7 + (row * 1.6)

        card = add_card_with_gradient_top(slide6, x, y, 4.0, 1.35, GOLD)
        card_frame = card.text_frame
        card_frame.text = f"{num}\n{title_text}\n{desc}"
        card_frame.paragraphs[0].font.size = Pt(32)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = PRIMARY
        card_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        card_frame.paragraphs[1].font.size = Pt(17)
        card_frame.paragraphs[1].font.bold = True
        card_frame.paragraphs[1].font.color.rgb = DARK
        card_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
        card_frame.paragraphs[2].font.size = Pt(11)
        card_frame.paragraphs[2].font.color.rgb = GRAY
        card_frame.paragraphs[2].alignment = PP_ALIGN.CENTER

    # ===== SLIDE 7: DATA SOURCES (NEW) =====
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    slide7.background.fill.solid()
    slide7.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame7 = title_box7.text_frame
    title_frame7.text = "Data Sources"
    title_p7 = title_frame7.paragraphs[0]
    title_p7.font.size = Pt(52)
    title_p7.font.bold = True
    title_p7.font.color.rgb = PRIMARY

    # Explanation
    explain_box7 = slide7.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.3), Inches(8), Inches(0.8)
    )
    explain_box7.fill.solid()
    explain_box7.fill.fore_color.rgb = PRIMARY
    explain_box7.line.fill.background()
    explain_frame7 = explain_box7.text_frame
    explain_frame7.text = "Multi-Source Data Aggregation\n\nWe integrate 8+ APIs to provide the most comprehensive DLMM analytics"
    explain_frame7.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in explain_frame7.paragraphs:
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    explain_frame7.paragraphs[0].font.size = Pt(22)
    explain_frame7.paragraphs[0].font.bold = True
    explain_frame7.paragraphs[2].font.size = Pt(13)

    # Data sources
    data_sources = [
        ("1", "Meteora API", "Pool data, positions, fees"),
        ("2", "DexScreener", "Price charts, volume stats"),
        ("3", "Jupiter", "Token prices, liquidity"),
        ("4", "RugCheck", "Security scans, risk analysis"),
        ("5", "Helius RPC", "Holder distribution, on-chain data"),
        ("6", "Birdeye", "Advanced charts, market data"),
        ("7", "Solscan", "Transaction explorer, token info"),
        ("8", "Solana RPC", "Real-time blockchain data")
    ]

    for i, (num, source, data) in enumerate(data_sources):
        row = i // 2
        col = i % 2
        x = 0.8 + (col * 4.4)
        y = 2.5 + (row * 1.15)

        card = add_card_with_gradient_top(slide7, x, y, 4.0, 1.0, SECONDARY)
        card_frame = card.text_frame
        card_frame.text = f"{num}  {source}\n{data}"
        card_frame.paragraphs[0].font.size = Pt(15)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = DARK
        card_frame.paragraphs[1].font.size = Pt(11)
        card_frame.paragraphs[1].font.color.rgb = GRAY

    # ===== SLIDE 8: TECH STACK =====
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    slide8.background.fill.solid()
    slide8.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame8 = title_box8.text_frame
    title_frame8.text = "Technology Stack"
    title_p8 = title_frame8.paragraphs[0]
    title_p8.font.size = Pt(52)
    title_p8.font.bold = True
    title_p8.font.color.rgb = PRIMARY

    tech = [
        "React 18", "Material-UI", "Flask", "Node.js",
        "PostgreSQL", "Supabase", "Solana Web3.js", "APScheduler",
        "Telegram Bot", "DexScreener", "Jupiter", "RugCheck"
    ]

    for i, tech_name in enumerate(tech):
        row = i // 4
        col = i % 4
        x = 1.0 + (col * 2.0)
        y = 1.8 + (row * 0.85)

        tech_box = slide8.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(1.8), Inches(0.65)
        )
        tech_box.fill.solid()
        tech_box.fill.fore_color.rgb = WHITE
        tech_box.line.color.rgb = PRIMARY
        tech_box.line.width = Pt(1.5)

        tech_frame = tech_box.text_frame
        tech_frame.text = tech_name
        tech_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tech_frame.paragraphs[0].font.size = Pt(12)
        tech_frame.paragraphs[0].font.bold = True
        tech_frame.paragraphs[0].font.color.rgb = DARK
        tech_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ===== SLIDE 8: COMPETITIVE ADVANTAGES =====
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    slide9.background.fill.solid()
    slide9.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame9 = title_box9.text_frame
    title_frame9.text = "Why We're Different"
    title_p9 = title_frame9.paragraphs[0]
    title_p9.font.size = Pt(52)
    title_p9.font.bold = True
    title_p9.font.color.rgb = PRIMARY

    advantages = [
        ("🎯", "100% DLMM Focused", "Built exclusively for Meteora"),
        ("🔒", "Security-First", "Integrated security analysis"),
        ("📲", "Mobile-First Alerts", "Always connected via Telegram"),
        ("⚙️", "Highly Customizable", "Adapts to your strategy")
    ]

    for i, (icon, title_text, desc) in enumerate(advantages):
        row = i // 2
        col = i % 2
        x = 0.8 + (col * 4.4)
        y = 1.7 + (row * 1.6)

        card = add_card_with_gradient_top(slide9, x, y, 4.0, 1.35, PRIMARY)
        card_frame = card.text_frame
        card_frame.text = f"{icon}  {title_text}\n{desc}"
        card_frame.paragraphs[0].font.size = Pt(17)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = DARK
        card_frame.paragraphs[1].font.size = Pt(12)
        card_frame.paragraphs[1].font.color.rgb = GRAY

    # ===== SLIDE 9: ROADMAP =====
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    slide10.background.fill.solid()
    slide10.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame10 = title_box10.text_frame
    title_frame10.text = "Product Roadmap"
    title_p10 = title_frame10.paragraphs[0]
    title_p10.font.size = Pt(52)
    title_p10.font.bold = True
    title_p10.font.color.rgb = PRIMARY

    roadmap = [
        ("Q1 2025", ["Automated trading", "Portfolio tracking", "Backtesting", "Watchlist"]),
        ("Q2 2025", ["Custom alerts", "Multi-DEX support", "Mobile apps", "API access"]),
        ("Q3 2025", ["Social features", "AI-powered scoring", "Cross-chain", "Enterprise"])
    ]

    for i, (quarter, items) in enumerate(roadmap):
        x = 1.0 + (i * 2.7)
        card = add_card_with_gradient_top(slide10, x, y=1.8, width=2.5, height=2.6, gradient_color=SECONDARY)
        card_frame = card.text_frame
        card_frame.text = f"{quarter}\n\n" + "\n".join(f"• {item}" for item in items)
        card_frame.paragraphs[0].font.size = Pt(20)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = PRIMARY
        card_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for j in range(2, len(card_frame.paragraphs)):
            card_frame.paragraphs[j].font.size = Pt(11)
            card_frame.paragraphs[j].font.color.rgb = GRAY

    # ===== SLIDE 10: BUSINESS MODEL =====
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    slide11.background.fill.solid()
    slide11.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box11 = slide11.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame11 = title_box11.text_frame
    title_frame11.text = "Business Model"
    title_p11 = title_frame11.paragraphs[0]
    title_p11.font.size = Pt(52)
    title_p11.font.bold = True
    title_p11.font.color.rgb = PRIMARY

    tiers = [
        ("Free", ["Core analytics", "Basic monitoring", "Security analysis", "Always free"]),
        ("Premium", ["$29/month", "1-min monitoring", "Unlimited alerts", "API access"]),
        ("Enterprise", ["Custom pricing", "White-label", "Dedicated infra", "SLA guarantees"])
    ]

    for i, (tier, features_list) in enumerate(tiers):
        x = 1.0 + (i * 2.7)
        card = add_card_with_gradient_top(slide11, x, y=1.8, width=2.5, height=2.6, gradient_color=GOLD)
        card_frame = card.text_frame
        card_frame.text = f"{tier}\n\n" + "\n".join(features_list)
        card_frame.vertical_anchor = MSO_ANCHOR.TOP
        card_frame.paragraphs[0].font.size = Pt(20)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = PRIMARY
        card_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for j in range(2, len(card_frame.paragraphs)):
            card_frame.paragraphs[j].font.size = Pt(11)
            card_frame.paragraphs[j].font.color.rgb = GRAY
            card_frame.paragraphs[j].alignment = PP_ALIGN.CENTER

    # ===== SLIDE 11: MARKET OPPORTUNITY =====
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    slide12.background.fill.solid()
    slide12.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box12 = slide12.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame12 = title_box12.text_frame
    title_frame12.text = "Market Opportunity"
    title_p12 = title_frame12.paragraphs[0]
    title_p12.font.size = Pt(52)
    title_p12.font.bold = True
    title_p12.font.color.rgb = PRIMARY

    market = [
        ("💰", "Growing DeFi Market", "Solana DeFi TVL exceeds $8B"),
        ("👥", "Target Audience", "Traders, LPs, protocols, funds"),
        ("📈", "Scalable Model", "Start Meteora, expand multi-chain"),
        ("🚀", "First-Mover", "No dedicated DLMM analytics")
    ]

    for i, (icon, title_text, desc) in enumerate(market):
        row = i // 2
        col = i % 2
        x = 0.8 + (col * 4.4)
        y = 1.7 + (row * 1.6)

        card = add_card_with_gradient_top(slide11, x, y, 4.0, 1.35, SECONDARY)
        card_frame = card.text_frame
        card_frame.text = f"{icon}  {title_text}\n{desc}"
        card_frame.paragraphs[0].font.size = Pt(17)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = DARK
        card_frame.paragraphs[1].font.size = Pt(12)
        card_frame.paragraphs[1].font.color.rgb = GRAY

    # ===== SLIDE 11: TRACTION =====
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    slide12.background.fill.solid()
    slide12.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box11 = slide12.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame11 = title_box11.text_frame
    title_frame11.text = "Current Traction"
    title_p11 = title_frame11.paragraphs[0]
    title_p11.font.size = Pt(52)
    title_p11.font.bold = True
    title_p11.font.color.rgb = PRIMARY

    stats11 = [("100%", "Feature Complete"), ("4K+", "Pools Tracked"), ("60s", "Refresh Rate")]
    for i, (num, label) in enumerate(stats11):
        x = 1.5 + (i * 2.5)
        card = add_card_with_gradient_top(slide12, x, y=1.8, width=2.3, height=1.6, gradient_color=PRIMARY)
        card_frame = card.text_frame
        card_frame.text = f"{num}\n{label}"
        card_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        card_frame.paragraphs[0].font.size = Pt(38)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = PRIMARY
        card_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        card_frame.paragraphs[1].font.size = Pt(15)
        card_frame.paragraphs[1].font.bold = True
        card_frame.paragraphs[1].font.color.rgb = DARK
        card_frame.paragraphs[1].alignment = PP_ALIGN.CENTER

    # Ready to launch box
    launch_box = slide12.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(4), Inches(6), Inches(1)
    )
    launch_box.fill.solid()
    launch_box.fill.fore_color.rgb = PRIMARY
    launch_box.line.fill.background()
    launch_frame = launch_box.text_frame
    launch_frame.text = "Ready to Launch\n\nPlatform is fully functional and tested"
    launch_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in launch_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    launch_frame.paragraphs[0].font.size = Pt(24)
    launch_frame.paragraphs[0].font.bold = True
    launch_frame.paragraphs[2].font.size = Pt(14)

    # ===== SLIDE 12: LIVE DEMO =====
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    slide13.background.fill.solid()
    slide13.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box12 = slide13.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame12 = title_box12.text_frame
    title_frame12.text = "Live Demo"
    title_p12 = title_frame12.paragraphs[0]
    title_p12.font.size = Pt(52)
    title_p12.font.bold = True
    title_p12.font.color.rgb = PRIMARY

    demo_box = slide13.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(2), Inches(7), Inches(1.2)
    )
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = PRIMARY
    demo_box.line.fill.background()
    demo_frame = demo_box.text_frame
    demo_frame.text = "Experience the Platform Live\n\nSee real-time data from 4,000+ pools"
    demo_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in demo_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    demo_frame.paragraphs[0].font.size = Pt(28)
    demo_frame.paragraphs[0].font.bold = True
    demo_frame.paragraphs[2].font.size = Pt(16)

    # CTA button
    cta = slide13.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(4), Inches(3), Inches(0.7)
    )
    cta.fill.solid()
    cta.fill.fore_color.rgb = GOLD
    cta.line.fill.background()
    cta_frame = cta.text_frame
    cta_frame.text = "Open Dashboard →"
    cta_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cta_frame.paragraphs[0].font.size = Pt(22)
    cta_frame.paragraphs[0].font.bold = True
    cta_frame.paragraphs[0].font.color.rgb = DARK
    cta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ===== SLIDE 13: CONTACT =====
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    slide14.background.fill.solid()
    slide14.background.fill.fore_color.rgb = LIGHT_GRAY

    title_box13 = slide14.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame13 = title_box13.text_frame
    title_frame13.text = "Let's Connect"
    title_p13 = title_frame13.paragraphs[0]
    title_p13.font.size = Pt(52)
    title_p13.font.bold = True
    title_p13.font.color.rgb = PRIMARY

    contacts = [("🌐", "Live Platform"), ("📧", "Email"), ("🐦", "Twitter"), ("💬", "Telegram")]
    for i, (icon, label) in enumerate(contacts):
        row = i // 2
        col = i % 2
        x = 1.5 + (col * 3.8)
        y = 1.6 + (row * 1.4)

        card = add_card_with_gradient_top(slide14, x, y, 3.4, 1.15, GOLD)
        card_frame = card.text_frame
        card_frame.text = f"{icon}  {label}\n[Your Info]"
        card_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        card_frame.paragraphs[0].font.size = Pt(16)
        card_frame.paragraphs[0].font.bold = True
        card_frame.paragraphs[0].font.color.rgb = DARK
        card_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        card_frame.paragraphs[1].font.size = Pt(12)
        card_frame.paragraphs[1].font.color.rgb = PRIMARY
        card_frame.paragraphs[1].alignment = PP_ALIGN.CENTER

    seeking_box = slide14.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(4.5), Inches(7), Inches(0.9)
    )
    seeking_box.fill.solid()
    seeking_box.fill.fore_color.rgb = PRIMARY
    seeking_box.line.fill.background()
    seeking_frame = seeking_box.text_frame
    seeking_frame.text = "We're Seeking\n\nStrategic Partners • Advisors • Beta Users • Feedback"
    seeking_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in seeking_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    seeking_frame.paragraphs[0].font.size = Pt(20)
    seeking_frame.paragraphs[0].font.bold = True
    seeking_frame.paragraphs[2].font.size = Pt(14)

    # ===== SLIDE 14: CLOSING HERO =====
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])
    slide15.background.fill.solid()
    slide15.background.fill.fore_color.rgb = PRIMARY

    # Add decorative shapes
    add_decorative_shapes(slide15, WHITE, SECONDARY)

    # Golden badge
    badge14 = slide15.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.8), Inches(1.5), Inches(4.4), Inches(0.5)
    )
    badge14.fill.solid()
    badge14.fill.fore_color.rgb = GOLD
    badge14.line.fill.background()
    badge14_frame = badge14.text_frame
    badge14_frame.text = "✨ THE FUTURE OF DLMM TRADING ✨"
    badge14_p = badge14_frame.paragraphs[0]
    badge14_p.font.size = Pt(14)
    badge14_p.font.bold = True
    badge14_p.font.color.rgb = DARK
    badge14_p.alignment = PP_ALIGN.CENTER
    badge14_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    title_box14 = slide15.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.3))
    title_frame14 = title_box14.text_frame
    title_frame14.text = "Ready to Trade Smarter?"
    title_p14 = title_frame14.paragraphs[0]
    title_p14.font.size = Pt(62)
    title_p14.font.bold = True
    title_p14.font.color.rgb = WHITE
    title_p14.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box14 = slide15.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(0.5))
    subtitle_frame14 = subtitle_box14.text_frame
    subtitle_frame14.text = "Join us in revolutionizing DLMM analytics on Solana"
    subtitle_p14 = subtitle_frame14.paragraphs[0]
    subtitle_p14.font.size = Pt(22)
    subtitle_p14.font.color.rgb = WHITE
    subtitle_p14.alignment = PP_ALIGN.CENTER

    # Thank you
    thanks_box = slide15.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank you for your time!\n\nQuestions? Let's discuss how we can work together."
    for p in thanks_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    thanks_frame.paragraphs[0].font.size = Pt(22)
    thanks_frame.paragraphs[0].font.bold = True
    thanks_frame.paragraphs[2].font.size = Pt(15)

    return prs

if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/younes/Documents/work/meteora-trending-pairs/meteora-trending-pairs/docs/Meteora_Trending_Pairs_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Beautiful PowerPoint created: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"✨ Design features: Rounded corners, shadows, gradient bars, decorative shapes")
